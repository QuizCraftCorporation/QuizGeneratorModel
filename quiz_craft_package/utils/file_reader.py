from langchain.document_loaders import PyPDFLoader, BSHTMLLoader, Docx2txtLoader
from pptx import Presentation

class FileReader:
    """
    Class to read different file extensions.
    """
    
    def __init__(self, filepath: str) -> None:
        """
        Read file using FileReader.

        Args:
            filepath (str): Path to file.
        """
        
        ext = filepath.split('.')[-1]
        if not ext in ['txt', 'pdf', 'html', 'docx', 'pptx']:
            raise Exception(f"Unsupported file format: {ext}")
        if ext == 'txt':
            input_file = open(filepath, "r", encoding="utf-8")
            self._content = input_file.read()
            input_file.close()
        elif ext == 'pdf':
            loader = PyPDFLoader(filepath)
            data = loader.load()
            self._content = ' '.join([page.page_content for page in data])
        elif ext == 'html':
            loader = BSHTMLLoader(filepath)
            data = loader.load()
            self._content = data[0].page_content
        elif ext == 'docx':
            loader = Docx2txtLoader(filepath)
            data = loader.load()
            self._content = data[0].page_content
        elif ext == 'pptx':
            data = ""
            data_file = Presentation(filepath)
            for slide in data_file.slides:
                for shapes in slide.shapes:
                    if shapes.has_text_frame:
                        data += shapes.text + "\n"
            self._content = data
        
    def get_content(self) -> str:
        """
        Extract content from a file.

        Returns:
            str: Text content of a file.
        """
        return self._content